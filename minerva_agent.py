import requests
import streamlit as st
import datetime as dt
import dashscope
from firecrawl import FirecrawlApp
import os
import json
from openai import OpenAI 
from apify_client import ApifyClient  # Added for X/Twitter scraping
import time
from datetime import datetime, timedelta
import csv
import pickle
import glob
import pandas as pd
import PyPDF2
import docx
from pptx import Presentation


# Create data directory if it doesn't exist
DATA_DIR = "data"
os.makedirs(DATA_DIR, exist_ok=True)

# File paths for persistent data storage
WEBSITE_DATA_PATH = os.path.join(DATA_DIR, "website_data.json")
WEBSITES_DB_PATH = os.path.join(DATA_DIR, "websites_db.json")
TWITTER_DATA_PATH = os.path.join(DATA_DIR, "twitter_data.json")
TWITTER_ACCOUNTS_DB_PATH = os.path.join(DATA_DIR, "twitter_accounts_db.json")
TWITTER_INSIGHTS_PATH = os.path.join(DATA_DIR, "twitter_insights.json")
RAG_DATA_PATH = os.path.join(DATA_DIR, "rag_data.pkl")
RAG_METADATA_PATH = os.path.join(DATA_DIR, "rag_metadata.json")

# Ensure forag directory exists
FORAG_DIR = "forag"
os.makedirs(FORAG_DIR, exist_ok=True)

os.environ["DASHSCOPE_API_KEY"] = "sk-1a28c3fcc7e044cbacd6faf47dc89755"

# Helper functions for websites database
def load_websites_db():
    """Load the websites database from JSON file"""
    if os.path.exists(WEBSITES_DB_PATH):
        with open(WEBSITES_DB_PATH, 'r', encoding='utf-8') as f:
            try:
                return json.load(f)
            except json.JSONDecodeError:
                return []
    else:
        # Initialize with default websites if database doesn't exist
        default_websites = [
            {"url": "www.qbitai.com", "name": "量子位"},
            {"url": "www.jiqizhixin.com", "name": "机器之心"},
            {"url": "36kr.com/information/AI", "name": "36kr"},
            {"url": "mp.sohu.com/profile?xpt=dXBob25lc3RjYXBpdGFsQDE2My5jb20=", "name": "硅兔赛跑"},
            {"url": "www.jazzyear.com", "name": "甲子光年"},
            {"url": "aiera.com.cn", "name": "新智元"},
            {"url": "wublock123.com", "name": "吴说Real"},
            {"url": "36kr.com/user/6038047", "name": "乌鸦智能说"},
            {"url": "www.zhihu.com/people/Khazix", "name": "数字生命卡兹克"}
        ]
        save_websites_db(default_websites)
        return default_websites

def save_websites_db(websites):
    """Save the websites database to a JSON file"""
    with open(WEBSITES_DB_PATH, 'w', encoding='utf-8') as f:
        json.dump(websites, f, ensure_ascii=False, indent=4)

def add_website_to_db(url, name):
    """Add a new website to the database"""
    websites = load_websites_db()
    
    # Check if website already exists
    for website in websites:
        if website["url"] == url:
            # Update the name if it's already in the database
            website["name"] = name
            save_websites_db(websites)
            return websites, False
    
    # Add new website
    websites.append({"url": url, "name": name})
    save_websites_db(websites)
    return websites, True

def remove_website_from_db(url):
    """Remove a website from the database"""
    websites = load_websites_db()
    websites = [website for website in websites if website["url"] != url]
    save_websites_db(websites)
    return websites

# Helper functions for Twitter accounts database
def load_twitter_accounts_db():
    """Load the Twitter accounts database from JSON file"""
    if os.path.exists(TWITTER_ACCOUNTS_DB_PATH):
        with open(TWITTER_ACCOUNTS_DB_PATH, 'r', encoding='utf-8') as f:
            try:
                return json.load(f)
            except json.JSONDecodeError:
                return []
    else:
        # Initialize with default Twitter accounts if database doesn't exist
        default_accounts = [
            {"handle": "sama", "name": "Sam Altman (OpenAI)"},
            {"handle": "ylecun", "name": "Yann LeCun (Meta)"},
            {"handle": "AndrewYNg", "name": "Andrew Ng (Landing AI)"},
            {"handle": "fchollet", "name": "François Chollet (Google)"},
            {"handle": "karpathy", "name": "Andrej Karpathy (Former Tesla)"},
            {"handle": "ilyasut", "name": "Ilya Sutskever (OpenAI)"}
        ]
        save_twitter_accounts_db(default_accounts)
        return default_accounts

def save_twitter_accounts_db(accounts):
    """Save the Twitter accounts database to a JSON file"""
    with open(TWITTER_ACCOUNTS_DB_PATH, 'w', encoding='utf-8') as f:
        json.dump(accounts, f, ensure_ascii=False, indent=4)

def add_twitter_account_to_db(handle, name):
    """Add a new Twitter account to the database"""
    accounts = load_twitter_accounts_db()
    
    # Clean handle by removing @ if present
    clean_handle = handle.replace("@", "").strip()
    
    # Check if account already exists
    for account in accounts:
        if account["handle"].lower() == clean_handle.lower():
            # Update the name if it's already in the database
            account["name"] = name
            save_twitter_accounts_db(accounts)
            return accounts, False
    
    # Add new account
    accounts.append({"handle": clean_handle, "name": name})
    save_twitter_accounts_db(accounts)
    return accounts, True

def remove_twitter_account_from_db(handle):
    """Remove a Twitter account from the database"""
    accounts = load_twitter_accounts_db()
    accounts = [account for account in accounts if account["handle"] != handle]
    save_twitter_accounts_db(accounts)
    return accounts

# Helper functions for data persistence
def save_website_data(domain, analysis):
    """Save website analysis data to a JSON file"""
    # Load existing data if available
    if os.path.exists(WEBSITE_DATA_PATH):
        with open(WEBSITE_DATA_PATH, 'r', encoding='utf-8') as f:
            try:
                data = json.load(f)
            except json.JSONDecodeError:
                data = {}
    else:
        data = {}
    
    # Add timestamp to track when the data was scraped
    data[domain] = {
        "analysis": analysis,
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    }
    
    # Save updated data
    with open(WEBSITE_DATA_PATH, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)
    
    return data

def load_website_data():
    """Load website analysis data from JSON file"""
    if os.path.exists(WEBSITE_DATA_PATH):
        with open(WEBSITE_DATA_PATH, 'r', encoding='utf-8') as f:
            try:
                return json.load(f)
            except json.JSONDecodeError:
                return {}
    return {}

def save_twitter_data(all_tweets, all_analyses):
    """Save Twitter data to JSON files"""
    # Save tweets and analyses
    twitter_data = {
        "tweets": all_tweets,
        "analyses": all_analyses,
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    }
    
    with open(TWITTER_DATA_PATH, 'w', encoding='utf-8') as f:
        json.dump(twitter_data, f, ensure_ascii=False, indent=4)
    
    return twitter_data

def save_twitter_insights(ai_insights, top_engaging_tweets):
    """Save Twitter insights to a JSON file"""
    insights_data = {
        "ai_insights": ai_insights,
        "top_engaging_tweets": top_engaging_tweets,
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    }
    
    with open(TWITTER_INSIGHTS_PATH, 'w', encoding='utf-8') as f:
        json.dump(insights_data, f, ensure_ascii=False, indent=4)
    
    return insights_data

def load_twitter_data():
    """Load Twitter data from JSON file"""
    if os.path.exists(TWITTER_DATA_PATH):
        with open(TWITTER_DATA_PATH, 'r', encoding='utf-8') as f:
            try:
                return json.load(f)
            except json.JSONDecodeError:
                return {"tweets": [], "analyses": [], "timestamp": None}
    return {"tweets": [], "analyses": [], "timestamp": None}

def load_twitter_insights():
    """Load Twitter insights from JSON file"""
    if os.path.exists(TWITTER_INSIGHTS_PATH):
        with open(TWITTER_INSIGHTS_PATH, 'r', encoding='utf-8') as f:
            try:
                return json.load(f)
            except json.JSONDecodeError:
                return {"ai_insights": None, "top_engaging_tweets": None, "timestamp": None}
    return {"ai_insights": None, "top_engaging_tweets": None, "timestamp": None}

def save_rag_data(local_facts, local_files):
    """Save RAG data using pickle for better handling of complex data structures"""
    rag_data = {
        "local_facts": local_facts,
        "local_files": local_files,
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    }
    
    with open(RAG_DATA_PATH, 'wb') as f:
        pickle.dump(rag_data, f)
    
    return rag_data

def load_rag_data():
    """Load RAG data from pickle file"""
    if os.path.exists(RAG_DATA_PATH):
        with open(RAG_DATA_PATH, 'rb') as f:
            try:
                return pickle.load(f)
            except (pickle.PickleError, EOFError):
                return {"local_facts": [], "local_files": [], "timestamp": None}
    return {"local_facts": [], "local_files": [], "timestamp": None}

# New functions for RAG metadata handling
def load_rag_metadata():
    """Load RAG metadata from JSON file"""
    if os.path.exists(RAG_METADATA_PATH):
        with open(RAG_METADATA_PATH, 'r', encoding='utf-8') as f:
            try:
                return json.load(f)
            except json.JSONDecodeError:
                return {}
    return {}

def save_rag_metadata(metadata):
    """Save RAG metadata to JSON file"""
    with open(RAG_METADATA_PATH, 'w', encoding='utf-8') as f:
        json.dump(metadata, f, ensure_ascii=False, indent=4)
    return metadata

def get_file_content(file_path):
    """Extracts content from a file based on its extension."""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    try:
        # Handle plain text and code files
        if ext in ['.txt', '.md', '.py', '.java', '.html', '.css', '.js', '.json']:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
        
        # Handle CSV files
        elif ext == '.csv':
            df = pd.read_csv(file_path)
            return df.to_string()
        
        # Handle Excel files
        elif ext in ['.xls', '.xlsx']:
            df = pd.read_excel(file_path)
            return df.to_string()
        
        # Handle PDF files
        elif ext == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                pages = []
                for page in reader.pages:
                    text = page.extract_text() or ""
                    pages.append(text)
                return "\n".join(pages)
        
        # Handle DOCX files
        elif ext == '.docx':
            doc = docx.Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs]
            return "\n".join(paragraphs)
        
        # Handle PPTX files
        elif ext == '.pptx':
            prs = Presentation(file_path)
            slide_texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_texts.append(shape.text)
            return "\n".join(slide_texts)
        
        # Fallback for other binary files: read the entire file in binary mode
        else:
            with open(file_path, 'rb') as f:
                content = f.read()
            return f"Binary file content ({ext}): {content}"
    
    except Exception as e:
        return f"Error reading file: {str(e)}"

def scan_forag_directory():
    """Scan the forag directory for files and update metadata"""
    metadata = load_rag_metadata()
    
    # Get all files in the forag directory
    files = glob.glob(os.path.join(FORAG_DIR, "*"))
    
    # Update metadata dict with default values for new files
    for file_path in files:
        file_name = os.path.basename(file_path)
        if file_name not in metadata:
            # Set default metadata for the file
            if file_name == "Sapient.inc 2025 Intro Booklet.pdf":
                # Special case for the sample file mentioned by the user
                metadata[file_name] = {
                    "类型": "项目",
                    "名称": "Sapient",
                    "标签": "宣传材料",
                    "content": get_file_content(file_path)
                }
            else:
                metadata[file_name] = {
                    "类型": "",
                    "名称": "",
                    "标签": "",
                    "content": get_file_content(file_path)
                }
    
    # Remove entries for files that no longer exist
    for file_name in list(metadata.keys()):
        if not os.path.exists(os.path.join(FORAG_DIR, file_name)):
            del metadata[file_name]
    
    # Save updated metadata
    save_rag_metadata(metadata)
    return metadata

# Function for chat using local factual knowledge (RAG)
def chat_with_local_facts(user_message, selected_filters=None):
    """Chat with local facts with optional filtering by metadata"""
    # Load metadata for all files
    metadata = load_rag_metadata()
    
    # Filter based on selected filters if provided
    filtered_metadata = metadata
    if selected_filters:
        filtered_metadata = {}
        for file_name, file_meta in metadata.items():
            # Check if the file matches all selected filters
            matches = True
            for key, value in selected_filters.items():
                if value and file_meta.get(key, "") != value:
                    matches = False
                    break
            if matches:
                filtered_metadata[file_name] = file_meta
    
    # Build context from filtered metadata
    context_text = ""
    for file_name, file_meta in filtered_metadata.items():
        # Add metadata information and content
        context_text += f"【文件】 {file_name}\n"
        context_text += f"类型: {file_meta.get('类型', '')}\n"
        context_text += f"名称: {file_meta.get('名称', '')}\n"
        context_text += f"标签: {file_meta.get('标签', '')}\n"
        context_text += f"内容: {file_meta.get('content', '')[:100000000]}\n\n"
    
    if not context_text:
        context_text = "当前没有匹配的本地文件信息。"
        
    messages = [
        {"role": "system", "content": f"你是一个基于本地事实知识库的智能助手。以下是部分文档内容用于辅助回答问题：\n{context_text}\n请基于这些内容回答用户问题。"},
        {"role": "user", "content": user_message},
    ]
    response = dashscope.Generation.call(
        api_key="sk-1a28c3fcc7e044cbacd6faf47dc89755",
        model="qwen-turbo",
        messages=messages,
        enable_search=True,
        result_format='message'
    )
    return response['output']['choices'][0]['message']['content']

# Initialize session state variables with persistent data if available
if "local_facts" not in st.session_state:
    rag_data = load_rag_data()
    st.session_state["local_facts"] = rag_data.get("local_facts", [])

if "local_files" not in st.session_state:
    if "local_facts" in st.session_state:  # Already loaded above
        rag_data = load_rag_data()
        st.session_state["local_files"] = rag_data.get("local_files", [])
    else:
        st.session_state["local_files"] = []

if "chat_history" not in st.session_state:
    st.session_state["chat_history"] = []

if "twitter_results" not in st.session_state:
    twitter_data = load_twitter_data()
    st.session_state["twitter_results"] = {
        "tweets": twitter_data.get("tweets", []),
        "analyses": twitter_data.get("analyses", [])
    }

if "ai_insights" not in st.session_state:
    insights_data = load_twitter_insights()
    st.session_state["ai_insights"] = insights_data.get("ai_insights")

if "top_engaging_tweets" not in st.session_state and "top_engaging_tweets" in load_twitter_insights():
    st.session_state["top_engaging_tweets"] = load_twitter_insights().get("top_engaging_tweets")

# Initialize the Firecrawl app with API key
fire_api = "fc-343fd362814545f295a89dc14ec4ee09"
app = FirecrawlApp(api_key=fire_api)
jina_api = "jina_26a656e516224ce28e71cc3b28fa7b07zUchXe4_MJ_935m8SpS9-TNGL--w"

# Function to get website name from URL
def get_website_name(url):
    websites = load_websites_db()
    
    # Clean URL by removing protocol and trailing slash
    clean_url = url.replace("https://", "").replace("http://", "").strip()
    if clean_url.endswith('/'):
        clean_url = clean_url[:-1]
    
    # Try to find the website in the database
    for website in websites:
        if website["url"] == clean_url or website["url"] in clean_url or clean_url in website["url"]:
            return website["name"]
    
    # If not found, return the URL itself
    return url

# Function to get raw HTML content from the websites using Firecrawl
def get_raw_html(domain):
    try:
        # Add protocol if not present
        if not domain.startswith('http'):
            url = f'https://{domain}'
        else:
            url = domain
            
        # Use Jina to get the raw HTML
        jina_url = f'https://r.jina.ai/{url}'
        headers = {
            'Authorization': 'Bearer jina_26a656e516224ce28e71cc3b28fa7b07zUchXe4_MJ_935m8SpS9-TNGL--w'
        }
        response = requests.get(jina_url, headers=headers)
        
        # Check if the request was successful
        response.raise_for_status()  # Raises an exception for 4XX/5XX responses
        
        # Return the raw HTML content
        return response.text
    except requests.exceptions.RequestException as e:
        return f"Error while fetching content from {domain}: {str(e)}"

# Function to prepare the message for Qwen LLM analysis
def analyze_with_qwen(domain, raw_html):
    messages = [
        {'role': 'system', 'content': 'You are a professional AI researcher. Analyze the raw HTML content and extract key topics in the following format: "1. Description | Website"'},
        {'role': 'user', 'content': f'''
        Analyze the raw HTML content from {domain} and provide the latest 10 topics with:
        1. Article titles in Chinese
        2. One-line descriptions in Chinese
        3. Website name
        Use current date: {dt.datetime.today().date()}.
        HTML Content: {raw_html}
        '''}
    ]
    response = dashscope.Generation.call(
        api_key="sk-1a28c3fcc7e044cbacd6faf47dc89755",
        model="qwen-turbo",
        messages=messages,
        enable_search=True,
        result_format='message'
    )
    return response['output']['choices'][0]['message']['content']

# Function to translate tweet text using Qwen
def translate_tweet_with_qwen(tweet_text):
    messages = [
        {'role': 'system', 'content': 'You are a professional translator. Translate the following text from English to Chinese.'},
        {'role': 'user', 'content': f'''
        Translate the following tweet from English to Chinese:
        
        {tweet_text}
        '''}
    ]
    
    response = dashscope.Generation.call(
        api_key="sk-1a28c3fcc7e044cbacd6faf47dc89755",
        model="qwen-turbo",
        messages=messages,
        result_format='message'
    )
    
    return response['output']['choices'][0]['message']['content']

# Function to batch translate tweets - MODIFIED to use Qwen for translation
def batch_translate_tweets(tweets):
    if not tweets:
        return tweets  # Return empty list if no tweets
    
    # Process each tweet individually
    for tweet in tweets:
        if 'text' in tweet and tweet['text'].strip():
            # Translate the text with Qwen
            translation = translate_tweet_with_qwen(tweet['text'])
            tweet['translation'] = translation
        else:
            tweet['translation'] = "无文本可翻译"
    
    return tweets

# Function to get top engaging tweets
def get_top_engaging_tweets(all_tweets):
    # Get top 5 by retweets, replies, and likes
    top_retweets = sorted(all_tweets, key=lambda x: x.get('retweets', 0), reverse=True)[:5]
    top_replies = sorted(all_tweets, key=lambda x: x.get('replies', 0), reverse=True)[:5]
    top_likes = sorted(all_tweets, key=lambda x: x.get('likes', 0), reverse=True)[:5]
    
    # Translate each category
    st.info("正在翻译热门推文...")
    translated_retweets = batch_translate_tweets(top_retweets)
    translated_replies = batch_translate_tweets(top_replies)
    translated_likes = batch_translate_tweets(top_likes)
    
    top_tweets = {
        'top_retweets': translated_retweets,
        'top_replies': translated_replies,
        'top_likes': translated_likes
    }
    
    # Save the top tweets data for persistence
    save_twitter_insights(st.session_state.get("ai_insights"), top_tweets)
    
    return top_tweets

# Function to scrape AI influencer tweets from X (Twitter)
def scrape_ai_influencer_tweets(twitter_handles):
    # Initialize the ApifyClient with your API token
    client = ApifyClient("apify_api_kbKxE4fYwbZOMBA30gS7DkbjinZqy91SEHb9")
    
    # Calculate date range for the last 2 days
    end_date = datetime.now().strftime("%Y-%m-%d")
    start_date = (datetime.now() - timedelta(days=2)).strftime("%Y-%m-%d")
    
    st.write(f"Scraping tweets from {start_date} to {end_date}")
    
    # Create directory for results if it doesn't exist
    os.makedirs("results", exist_ok=True)
    
    # Initialize results storage
    all_tweets = []
    all_analyses = []
    
    # Progress tracking for Streamlit
    progress = st.progress(0)
    status_text = st.empty()
    
    # CSV for structured data
    with open("results/ai_influencer_tweets.csv", "w", newline="", encoding="utf-8") as csvfile:
        csv_writer = csv.writer(csvfile)
        # Write header
        csv_writer.writerow(["Author", "Twitter Handle", "text", "createdAt", "replyCount", "likeCount", "retweetCount", "url"])
        
        # Load the Twitter accounts database to get names
        twitter_accounts = load_twitter_accounts_db()
        account_names = {account["handle"]: account["name"] for account in twitter_accounts}
        
        # Iterate through each Twitter handle
        for i, handle in enumerate(twitter_handles):
            status_text.text(f"[{i+1}/{len(twitter_handles)}] Scraping tweets for @{handle}...")
            progress.progress((i+1) / len(twitter_handles))
            
            # Get display name from database or use handle as fallback
            display_name = account_names.get(handle, handle)
            
            # Prepare the Actor input
            run_input = {
                "maxItems": 100,  # Limit to 100 tweets per author
                "sort": "Latest",
                "author": handle,
                "start": start_date,
                "end": end_date,
            }
            
            # Implement retry logic
            max_retries = 3
            retries = 0
            success = False
            
            while retries < max_retries and not success:
                try:
                    # Run the Actor and wait for it to finish
                    run = client.actor("61RPP7dywgiy0JPD0").call(run_input=run_input)
                    
                    # Track tweets for this author
                    author_tweets = []
                    
                    # Fetch tweets from the dataset
                    for tweet in client.dataset(run["defaultDatasetId"]).iterate_items():
                        # Extract relevant data
                        author_name = tweet.get("authorName", "")
                        text = tweet.get("text", "")
                        date = tweet.get("createdAt", "")
                        replies = tweet.get("replyCount", 0)
                        likes = tweet.get("likeCount", 0)
                        retweets = tweet.get("retweetCount", 0)
                        url = tweet.get("url", "")
                        
                        # Create tweet object
                        tweet_data = {
                            "author": author_name,
                            "handle": handle,
                            "display_name": display_name,  # Add display name from database
                            "text": text,
                            "date": date,
                            "replies": replies,
                            "likes": likes,
                            "retweets": retweets,
                            "url": url
                        }
                        
                        # Add to collections
                        author_tweets.append(tweet_data)
                        all_tweets.append(tweet_data)
                        
                        # Write to CSV
                        csv_writer.writerow([author_name, handle, text, date, replies, likes, retweets, url])
                    
                    # Save individual author results to JSON
                    if author_tweets:
                        with open(f"results/{handle}_tweets.json", "w", encoding="utf-8") as f:
                            json.dump(author_tweets, f, indent=4)
                        
                        # Analyze tweets with Qwen
                        analysis = analyze_tweets_with_qwen(handle, author_tweets)
                        # Get the author's full name from the database or first tweet
                        author_full_name = display_name if display_name != handle else (author_tweets[0]["author"] if author_tweets else handle)
                        all_analyses.append({
                            "handle": handle,
                            "author_name": author_full_name,
                            "analysis": analysis
                        })
                    
                    st.write(f"  Scraped {len(author_tweets)} tweets for @{handle} ({display_name})")
                    success = True
                    
                except Exception as e:
                    retries += 1
                    st.error(f"  Error (attempt {retries}/{max_retries}): {e}")
                    if retries < max_retries:
                        st.info(f"  Retrying in 10 seconds...")
                        time.sleep(10)
                    else:
                        st.error(f"  Failed to scrape tweets for @{handle} after {max_retries} attempts.")
            
            # Add a delay between authors to respect rate limits
            if i < len(twitter_handles) - 1:  # Don't delay after the last author
                time.sleep(3)
    
    # After collecting all tweets, analyze them collectively for AI insights
    if all_tweets:
        # Changed to use Deepseek for AI insights
        ai_insights = extract_ai_insights_with_deepseek(all_tweets)
        st.session_state["ai_insights"] = ai_insights
        
        # Get and translate top engaging tweets
        st.info("正在获取和翻译最具互动性的推文...")
        top_tweets = get_top_engaging_tweets(all_tweets)
        st.session_state["top_engaging_tweets"] = top_tweets
        
        # Save all Twitter data for persistence
        save_twitter_data(all_tweets, all_analyses)
        save_twitter_insights(ai_insights, top_tweets)
    
    # Final summary
    progress.empty()
    status_text.empty()
    st.success(f"Scraped a total of {len(all_tweets)} tweets from {len(twitter_handles)} AI influencers")
    
    return all_tweets, all_analyses

# Function to analyze tweets with Qwen
def analyze_tweets_with_qwen(handle, tweets_data):
    # Prepare tweet text for analysis
    tweet_content = ""
    for tweet in tweets_data:
        tweet_content += f"Tweet: {tweet['text']}\nDate: {tweet['date']}\nURL: {tweet['url']}\n\n"
    
    messages = [
        {'role': 'system', 'content': 'You are a professional AI researcher. Analyze the tweets and extract key insights and topics.'},
        {'role': 'user', 'content': f'''
        Analyze the following tweets from AI influencer @{handle} and provide in Chinese:
        1. key topics or trends in Chinese
        2. Brief summary of main points in Chinese
        3. Any significant announcements or news in Chinese
        
        Current date: {datetime.now().strftime('%Y-%m-%d')}.
        
        Tweets:
        {tweet_content}
        '''}
    ]
    
    response = dashscope.Generation.call(
        api_key="sk-1a28c3fcc7e044cbacd6faf47dc89755",
        model="qwen-turbo",
        messages=messages,
        enable_search=True,
        result_format='message'
    )
    
    return response['output']['choices'][0]['message']['content']

# New function to extract AI insights from collected tweets using Deepseek
def extract_ai_insights_with_deepseek(tweets_data):
    # Organize tweets by their importance (using likes, retweets as indicators)
    sorted_tweets = sorted(
        tweets_data, 
        key=lambda x: (x.get('likes', 0) + x.get('retweets', 0)), 
        reverse=True
    )
    
    # Take top 100 tweets or all if less than 100
    top_tweets = sorted_tweets[:min(100, len(sorted_tweets))]
    
    # Prepare tweet content for analysis
    tweet_content = ""
    for tweet in top_tweets:
        tweet_content += f"Author: @{tweet['handle']} ({tweet.get('display_name', tweet['author'])})\nTweet: {tweet['text']}\nLikes: {tweet['likes']}, Retweets: {tweet['retweets']}\nDate: {tweet['date']}\n\n"
    
    # Use OpenAI client with Deepseek model
    client = OpenAI(
        api_key= "sk-1a28c3fcc7e044cbacd6faf47dc89755",
        base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
    )
    
    messages = [
        {"role": "system", "content": "You are an expert AI researcher and analyst. Your task is to identify emerging trends, groundbreaking research, and industry shifts in artificial intelligence by analyzing tweets from leading AI professionals."},
        {"role": "user", "content": f'''
        Analyze the following collection of tweets from leading AI professionals and provide a comprehensive analysis of his/her AI insights:
        1. 最新AI技术趋势 (Latest AI Technology Trends): Identify emerging technologies or approaches that appear to be gaining momentum in the AI community.
        
        2. 研究方向前沿 (Research Frontiers): Extract research areas that seem to be at the cutting edge based on researcher discussions.
        
        3. 行业发展动态 (Industry Developments): Summarize key business or industry shifts that are evident from these tweets.
                
        4. 未来AI展望 (Future AI Outlook): Based on these experts' tweets, provide insights on where AI might be heading in the near future.
        
        Current date: {datetime.now().strftime('%Y-%m-%d')}.
        
        Tweets:
        {tweet_content}
        
        Provide your analysis in Chinese, with clear section headers and bullet points where appropriate, do not use Markdown format.
        '''}
    ]
    
    completion = client.chat.completions.create(
        model="deepseek-r1",
        messages=messages
    )
    
    # 获取思考过程和最终答案
    thinking_process = completion.choices[0].message.reasoning_content
    ai_insights = completion.choices[0].message.content
    
    # Save AI insights for persistence
    save_twitter_insights(ai_insights, st.session_state.get("top_engaging_tweets"))
    
    return {
        "思考过程": thinking_process,
        "分析结果": ai_insights
    }

# Function for direct chat using Qwen (standard mode)
def chat_with_qwen(user_message):
    messages = [
        {"role": "system", "content": "You are a helpful assistant."},
        {"role": "user", "content": user_message},
    ]
    response = dashscope.Generation.call(
        api_key="sk-1a28c3fcc7e044cbacd6faf47dc89755",
        model="qwen-turbo",
        messages=messages,
        enable_search=True,
        result_format='message'
    )
    return response['output']['choices'][0]['message']['content']

# Function for direct chat using Deepseek model
def chat_with_deepseek(user_message):
    messages = [
        {"role": "system", "content": "You are a helpful assistant."},
        {"role": "user", "content": user_message},
    ]
    completion = client.chat.completions.create(
        model="deepseek-r1",
        messages=messages
    )
    
    # 通过reasoning_content字段获取思考过程
    thinking_process = completion.choices[0].message.reasoning_content
    final_answer = completion.choices[0].message.content
    
    return {
        "思考过程": thinking_process,
        "最终答案": final_answer
    }

# ----------------------- Streamlit UI -----------------------
st.title("Minerva Agent")
# Create tabs for different functionalities
tabs = st.tabs(["热点监控", "事实知识库 (RAG)", "直接聊天"])

# ----------------------- Tab 1: Trending Topics Monitoring -----------------------
with tabs[0]:
    st.header("热点监控")
    
    # Create subtabs for different monitoring sources
    monitoring_tabs = st.tabs(["网站监控", "X/Twitter监控"])
    
    # Website monitoring tab
    with monitoring_tabs[0]:
        st.write("监控推流的各大信息网站的热点")
        
        # Load existing website data
        website_data = load_website_data()
        if website_data:
            st.info(f"已加载 {len(website_data)} 个网站的历史数据。上次更新时间: {list(website_data.values())[0].get('timestamp', '未知')}")
        
        # Load websites from database
        websites_db = load_websites_db()
        
        # Website database management section
        with st.expander("网站数据库管理"):
            # Display current websites in the database
            st.subheader("当前监控的网站")
            for i, website in enumerate(websites_db):
                col1, col2, col3 = st.columns([3, 3, 1])
                with col1:
                    st.text(website["name"])
                with col2:
                    st.text(website["url"])
                with col3:
                    if st.button("删除", key=f"del_website_{i}"):
                        websites_db = remove_website_from_db(website["url"])
                        st.success(f"已删除网站: {website['name']}")
                        st.rerun()
            
            # Form to add a new website
            st.subheader("添加新网站")
            with st.form("add_website_form"):
                new_url = st.text_input("网站URL (例如: www.example.com):")
                new_name = st.text_input("网站名称 (例如: 示例网站):")
                submitted = st.form_submit_button("添加网站")
                
                if submitted and new_url and new_name:
                    # Clean URL
                    clean_url = new_url.replace("https://", "").replace("http://", "").strip()
                    if clean_url.endswith('/'):
                        clean_url = clean_url[:-1]
                    
                    websites_db, is_new = add_website_to_db(clean_url, new_name)
                    if is_new:
                        st.success(f"已添加新网站: {new_name} ({clean_url})")
                    else:
                        st.success(f"已更新网站名称: {new_name} ({clean_url})")
                    st.rerun()
        
        # Build a list of websites for monitoring
        monitoring_websites = []
        selected_websites = st.multiselect(
            "选择要监控的网站:",
            options=[website["name"] for website in websites_db],
            default=[website["name"] for website in websites_db]
        )
        
        for website in websites_db:
            if website["name"] in selected_websites:
                monitoring_websites.append(website["url"])
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("开始网站监控", key="scrape_websites"):
                for site in monitoring_websites:
                    # Get the website name
                    site_name = get_website_name(site)
                    st.write(f"### 正在拉取 {site_name} 的数据...")
                    raw_html = get_raw_html(site)
                    if isinstance(raw_html, str) and ('Error' in raw_html or 'Failed' in raw_html):
                        st.error(raw_html)
                    else:
                        st.write("数据拉取成功，正在分析热点内容...")
                        analysis = analyze_with_qwen(site, raw_html)
                        # Save the analysis for this website
                        save_website_data(site, analysis)
                        st.text_area(f"{site_name} 热点分析", analysis, height=300)
                    st.markdown("---")
        
        with col2:
            if st.button("重新抓取", key="rescrape_websites"):
                st.warning("正在重新抓取所有网站数据...")
                # Clear existing data
                if os.path.exists(WEBSITE_DATA_PATH):
                    os.remove(WEBSITE_DATA_PATH)
                
                # Perform scraping as in the "开始网站监控" button
                for site in monitoring_websites:
                    # Get the website name
                    site_name = get_website_name(site)
                    st.write(f"### 正在重新拉取 {site_name} 的数据...")
                    raw_html = get_raw_html(site)
                    if isinstance(raw_html, str) and ('Error' in raw_html or 'Failed' in raw_html):
                        st.error(raw_html)
                    else:
                        st.write("数据拉取成功，正在分析热点内容...")
                        analysis = analyze_with_qwen(site, raw_html)
                        # Save the analysis for this website
                        save_website_data(site, analysis)
                        st.text_area(f"{site_name} 热点分析", analysis, height=300)
                    st.markdown("---")
        
        # Display cached data if not scraping
        hide_web_data = st.button("隐藏历史数据", key="hide_website_data")
        if website_data and not hide_web_data:
            st.subheader("已加载的网站数据")
            for site, data in website_data.items():
                # Get the website name for display
                site_name = get_website_name(site)
                with st.expander(f"{site_name} - 上次更新: {data.get('timestamp', '未知')}"):
                    st.text_area(f"{site_name} 缓存热点分析", data.get('analysis', '没有数据'), height=300)
    
    # X/Twitter monitoring tab
    with monitoring_tabs[1]:
        st.write("监控AI领域专家X动态")
        
        # Display information about the scraper
        st.info("这个功能会抓取AI领域专家的X动态，并用Deepseek进行分析，提取insights。")
        
        # Load existing Twitter data
        twitter_data = load_twitter_data()
        twitter_insights = load_twitter_insights()
        
        if twitter_data.get("tweets"):
            st.info(f"已加载 {len(twitter_data.get('tweets', []))} 条推文数据。上次更新时间: {twitter_data.get('timestamp', '未知')}")
        
        # Load Twitter accounts from database
        twitter_accounts_db = load_twitter_accounts_db()
        
        # Twitter accounts database management section
        with st.expander("X/Twitter账号数据库管理"):
            # Display current Twitter accounts in the database
            st.subheader("当前监控的X/Twitter账号")
            for i, account in enumerate(twitter_accounts_db):
                col1, col2, col3 = st.columns([3, 3, 1])
                with col1:
                    st.text(account["name"])
                with col2:
                    st.text(f"@{account['handle']}")
                with col3:
                    if st.button("删除", key=f"del_account_{i}"):
                        twitter_accounts_db = remove_twitter_account_from_db(account["handle"])
                        st.success(f"已删除账号: {account['name']} (@{account['handle']})")
                        st.rerun()
            
            # Form to add a new Twitter account
            st.subheader("添加新X/Twitter账号")
            with st.form("add_twitter_account_form"):
                new_handle = st.text_input("账号名称 (例如: elonmusk):")
                new_name = st.text_input("显示名称 (例如: Elon Musk (Tesla/X)):")
                submitted = st.form_submit_button("添加账号")
                
                if submitted and new_handle and new_name:
                    # Clean handle by removing @ if present
                    clean_handle = new_handle.replace("@", "").strip()
                    
                    twitter_accounts_db, is_new = add_twitter_account_to_db(clean_handle, new_name)
                    if is_new:
                        st.success(f"已添加新账号: {new_name} (@{clean_handle})")
                    else:
                        st.success(f"已更新账号名称: {new_name} (@{clean_handle})")
                    st.rerun()
        
        # Build a list of Twitter handles for monitoring
        monitoring_handles = []
        selected_accounts = st.multiselect(
            "选择要监控的X/Twitter账号:",
            options=[f"{account['name']} (@{account['handle']})" for account in twitter_accounts_db],
            default=[f"{account['name']} (@{account['handle']})" for account in twitter_accounts_db[:3]]
        )
        
        for account_display in selected_accounts:
            # Extract handle from display string like "Sam Altman (@sama)"
            handle = account_display.split("(@")[1].replace(")", "")
            monitoring_handles.append(handle)
        
        # Limit the number of selected handles to avoid rate limiting
        max_handles = min(len(monitoring_handles) if monitoring_handles else 3, 50000)
        if len(monitoring_handles) > max_handles:
            st.warning(f"由于API限制，一次最多只能监控{max_handles}个账号。已自动选择前{max_handles}个账号。")
            monitoring_handles = monitoring_handles[:max_handles]
        
        col1, col2 = st.columns(2)
        with col1:
            # Scrape button
            if st.button("开始抓取X数据", key="scrape_twitter"):
                if monitoring_handles:
                    st.write(f"### 正在抓取 {len(monitoring_handles)} 个AI专家的Twitter数据...")
                    
                    # Call the function to scrape and analyze tweets
                    all_tweets, all_analyses = scrape_ai_influencer_tweets(monitoring_handles)
                    
                    # Store in session state for persistence
                    st.session_state["twitter_results"] = {
                        "tweets": all_tweets,
                        "analyses": all_analyses
                    }
                    
                    # Display analyses
                    if all_analyses:
                        # First display the collective AI insights
                        if st.session_state["ai_insights"]:
                            st.subheader("🔍 AI行业综合洞察")
                            st.text_area("AI行业洞察分析", st.session_state["ai_insights"], height=400)
                        
                        # Display top engaging tweets with translations
                        if "top_engaging_tweets" in st.session_state:
                            st.subheader("🔝 最具互动性的推文")
                            
                            top_tweets_tabs = st.tabs(["热门转发", "热门回复", "热门点赞"])
                            
                            # Display top retweets
                            with top_tweets_tabs[0]:
                                st.write("### 最高转发量推文")
                                for i, tweet in enumerate(st.session_state["top_engaging_tweets"]["top_retweets"], 1):
                                    with st.expander(f"{i}. @{tweet['handle']} (转发: {tweet['retweets']})"):
                                        # Use display_name if available
                                        display_name = tweet.get('display_name', tweet['author'])
                                        st.markdown(f"""
                                        **作者：** {display_name} (@{tweet['handle']})  
                                        **日期：** {tweet['date']}  
                                        **原文：** {tweet['text']}  
                                        **中文翻译：** {tweet['translation']}  
                                        **互动：** 👍 {tweet['likes']} | 🔁 {tweet['retweets']} | 💬 {tweet['replies']}  
                                        **链接：** [查看原文]({tweet['url']})
                                        """)
                            
                            # Display top replies
                            with top_tweets_tabs[1]:
                                st.write("### 最高回复量推文")
                                for i, tweet in enumerate(st.session_state["top_engaging_tweets"]["top_replies"], 1):
                                    with st.expander(f"{i}. @{tweet['handle']} (回复: {tweet['replies']})"):
                                        # Use display_name if available
                                        display_name = tweet.get('display_name', tweet['author'])
                                        st.markdown(f"""
                                        **作者：** {display_name} (@{tweet['handle']})  
                                        **日期：** {tweet['date']}  
                                        **原文：** {tweet['text']}  
                                        **中文翻译：** {tweet['translation']}  
                                        **互动：** 👍 {tweet['likes']} | 🔁 {tweet['retweets']} | 💬 {tweet['replies']}  
                                        **链接：** [查看原文]({tweet['url']})
                                        """)
                            
                            # Display top likes
                            with top_tweets_tabs[2]:
                                st.write("### 最高点赞量推文")
                                for i, tweet in enumerate(st.session_state["top_engaging_tweets"]["top_likes"], 1):
                                    with st.expander(f"{i}. @{tweet['handle']} (点赞: {tweet['likes']})"):
                                        # Use display_name if available
                                        display_name = tweet.get('display_name', tweet['author'])
                                        st.markdown(f"""
                                        **作者：** {display_name} (@{tweet['handle']})  
                                        **日期：** {tweet['date']}  
                                        **原文：** {tweet['text']}  
                                        **中文翻译：** {tweet['translation']}  
                                        **互动：** 👍 {tweet['likes']} | 🔁 {tweet['retweets']} | 💬 {tweet['replies']}  
                                        **链接：** [查看原文]({tweet['url']})
                                        """)
                        
                        # Then display individual analyses
                        st.subheader("个人推文分析")
                        for analysis_item in all_analyses:
                            handle = analysis_item["handle"]
                            author_name = analysis_item.get("author_name", handle)
                            analysis = analysis_item["analysis"]
                            
                            # Get tweets for this handle
                            handle_tweets = [t for t in all_tweets if t["handle"] == handle]
                            
                            # Create a container for each author
                            author_container = st.container()
                            with author_container:
                                st.markdown(f"### {author_name} (@{handle})")
                                
                                # Display analysis in the container
                                st.text_area(f"{handle} 推文分析", analysis, height=250)
                                
                                # Display tweets for this author
                                if handle_tweets:
                                    # Display raw tweets in a simple list format to avoid nesting expanders
                                    st.markdown(f"#### @{handle} 的原始推文 ({len(handle_tweets)} 条)")
                                    for tweet in handle_tweets:
                                        st.markdown(f"""
                                        **日期：** {tweet['date']}  
                                        **内容：** {tweet['text']}  
                                        **互动：** 👍 {tweet['likes']} | 🔁 {tweet['retweets']} | 💬 {tweet['replies']}  
                                        **链接：** [查看原文]({tweet['url']})
                                        ---
                                        """)
                                else:
                                    st.info(f"未找到 @{handle} 的推文数据")
                else:
                    st.warning("请选择至少一个X账号进行监控。")
        
        with col2:
            if st.button("重新抓取", key="rescrape_twitter"):
                if monitoring_handles:
                    st.warning("正在重新抓取所有X数据...")
                    # Clear existing data
                    if os.path.exists(TWITTER_DATA_PATH):
                        os.remove(TWITTER_DATA_PATH)
                    if os.path.exists(TWITTER_INSIGHTS_PATH):
                        os.remove(TWITTER_INSIGHTS_PATH)
                    
                    st.write(f"### 正在重新抓取这2天内 {len(monitoring_handles)} 个AI专家的Twitter数据，并用Deepseek分析（会有些慢）...")
                    
                    # Call the function to scrape and analyze tweets
                    all_tweets, all_analyses = scrape_ai_influencer_tweets(monitoring_handles)
                    
                    # Store in session state for persistence
                    st.session_state["twitter_results"] = {
                        "tweets": all_tweets,
                        "analyses": all_analyses
                    }
                    
                    # Similar display code as above
                else:
                    st.warning("请选择至少一个X账号进行监控。")
        
        # Display cached Twitter data if available and not scraping
        hide_twitter_data = st.button("隐藏历史数据", key="hide_twitter_data")
        if (twitter_data.get("tweets") or twitter_insights.get("ai_insights")) and not hide_twitter_data:
            st.subheader("已加载的X数据")
            
            # Display AI insights
            if twitter_insights.get("ai_insights"):
                with st.expander(f"AI行业综合洞察 - 上次更新: {twitter_insights.get('timestamp', '未知')}"):
                    st.text_area("缓存AI行业洞察分析", twitter_insights.get("ai_insights", "没有数据"), height=400)
            
            # Display top engaging tweets if available
            if twitter_insights.get("top_engaging_tweets"):
                with st.expander("最具互动性的推文"):
                    top_engaging_tweets = twitter_insights.get("top_engaging_tweets", {})
                    
                    if "top_retweets" in top_engaging_tweets:
                        st.markdown("#### 最高转发量推文")
                        for i, tweet in enumerate(top_engaging_tweets["top_retweets"], 1):
                            display_name = tweet.get('display_name', tweet['author'])
                            st.markdown(f"""
                            **{i}. {display_name} (@{tweet['handle']}) (转发: {tweet['retweets']})** - {tweet['text']}
                            """)
                    
                    if "top_likes" in top_engaging_tweets:
                        st.markdown("#### 最高点赞量推文")
                        for i, tweet in enumerate(top_engaging_tweets["top_likes"], 1):
                            display_name = tweet.get('display_name', tweet['author'])
                            st.markdown(f"""
                            **{i}. {display_name} (@{tweet['handle']}) (点赞: {tweet['likes']})** - {tweet['text']}
                            """)
                    if "top_replies" in top_engaging_tweets:
                        st.markdown("#### 最高回复量推文")
                        for i, tweet in enumerate(top_engaging_tweets["top_replies"], 1):
                            display_name = tweet.get('display_name', tweet['author'])
                            st.markdown(f"""
                            **{i}. {display_name} (@{tweet['handle']}) (回复: {tweet['replies']})** - {tweet['text']}
                            """)
            
            # Display individual analyses if available
            if twitter_data.get("analyses"):
                with st.expander("个人推文分析"):
                    for analysis_item in twitter_data.get("analyses", []):
                        handle = analysis_item["handle"]
                        author_name = analysis_item.get("author_name", handle)
                        st.markdown(f"**{author_name} (@{handle})**")
                        st.markdown(analysis_item["analysis"])
                        st.markdown("---")

# ----------------------- Tab 2: RAG (Knowledge Base) -----------------------
with tabs[1]:
    st.header("事实知识库 (RAG)")
    
    # Create two subtabs - one for file management, one for chatting
    rag_tabs = st.tabs(["文件管理", "基于知识库聊天"])
    
    # File management tab
    with rag_tabs[0]:
        st.write("管理本地文件知识库")
        
        # Scan the forag directory for files and update metadata
        metadata = scan_forag_directory()
        
        # Display current files and metadata
        st.subheader("当前文件库")
        
        if not metadata:
            st.info("目前没有发现任何文件。请上传文件或将文件放入 'forag' 文件夹。")
        else:
            # Display files in a table format
            for i, (file_name, file_meta) in enumerate(metadata.items()):
                with st.expander(f"{i+1}. {file_name}"):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.write("**当前元数据：**")
                        st.write(f"**类型:** {file_meta.get('类型', '未设置')}")
                        st.write(f"**名称:** {file_meta.get('名称', '未设置')}")
                        st.write(f"**标签:** {file_meta.get('标签', '未设置')}")
                    
                    with col2:
                        # Form to update metadata
                        with st.form(key=f"update_metadata_{i}"):
                            st.write("**更新元数据：**")
                            new_type = st.text_input("类型", value=file_meta.get('类型', ''), key=f"type_{i}")
                            new_name = st.text_input("名称", value=file_meta.get('名称', ''), key=f"name_{i}")
                            new_tag = st.text_input("标签", value=file_meta.get('标签', ''), key=f"tag_{i}")
                            
                            if st.form_submit_button("更新元数据"):
                                metadata[file_name]['类型'] = new_type
                                metadata[file_name]['名称'] = new_name
                                metadata[file_name]['标签'] = new_tag
                                save_rag_metadata(metadata)
                                st.success(f"已更新 {file_name} 的元数据")
                                st.rerun()
                    
                    # Display a preview of the file content
                    content_preview = file_meta.get('content', '')
                    if content_preview:
                        if len(content_preview) > 100000000:
                            content_preview = content_preview[:100000000] + "..."
                        st.write("**文件内容预览：**")
                        st.text_area("内容", content_preview, height=200, key=f"content_{i}")
                    
                    # Button to delete file from metadata (doesn't delete the actual file)
                    if st.button("从知识库中移除", key=f"remove_file_{i}"):
                        del metadata[file_name]
                        save_rag_metadata(metadata)
                        st.success(f"已从知识库中移除 {file_name}")
                        st.rerun()
        
        # Upload new files section
        st.subheader("上传新文件")
        with st.form("upload_new_file_form", clear_on_submit=True):
            uploaded_file = st.file_uploader("选择要上传的文件", key="rag_uploader")
            
            # Metadata fields
            file_type = st.text_input("类型 (例如: 项目, 产品, 文章)")
            file_name = st.text_input("名称 (例如: Sapient)")
            file_tag = st.text_input("标签 (例如: 宣传材料, 技术文档)")
            
            if st.form_submit_button("上传文件"):
                if uploaded_file is not None:
                    # Save the file to the forag directory
                    file_path = os.path.join(FORAG_DIR, uploaded_file.name)
                    with open(file_path, "wb") as f:
                        f.write(uploaded_file.getbuffer())
                    
                    # Extract content
                    file_content = get_file_content(file_path)
                    
                    # Update metadata
                    metadata = load_rag_metadata()
                    metadata[uploaded_file.name] = {
                        "类型": file_type,
                        "名称": file_name,
                        "标签": file_tag,
                        "content": file_content
                    }
                    save_rag_metadata(metadata)
                    
                    st.success(f"文件 {uploaded_file.name} 已上传并添加到知识库")
                    st.rerun()
                else:
                    st.error("请选择一个文件上传")
        
        # Refresh button
        if st.button("刷新文件库"):
            scan_forag_directory()
            st.rerun()
        
        # Reset button - clears all metadata
        if st.button("重置知识库元数据"):
            save_rag_metadata({})
            st.success("知识库元数据已重置")
            st.rerun()
    
    # Chat with knowledge base tab
    with rag_tabs[1]:
        st.write("与本地知识库对话")
        
        # Load metadata
        metadata = load_rag_metadata()
        
        if not metadata:
            st.warning('知识库中没有文件。请先在"文件管理"选项卡中添加文件。')
        else:
            # Get distinct types, names, and tags for filtering
            all_types = list(set(file_meta.get('类型', '') for file_meta in metadata.values() if file_meta.get('类型', '')))
            all_names = list(set(file_meta.get('名称', '') for file_meta in metadata.values() if file_meta.get('名称', '')))
            all_tags = list(set(file_meta.get('标签', '') for file_meta in metadata.values() if file_meta.get('标签', '')))
            
            # Add empty option
            all_types = [""] + all_types
            all_names = [""] + all_names
            all_tags = [""] + all_tags
            
            # Filters
            st.subheader("过滤选项")
            col1, col2, col3 = st.columns(3)
            with col1:
                selected_type = st.selectbox("类型", all_types)
            with col2:
                selected_name = st.selectbox("名称", all_names)
            with col3:
                selected_tag = st.selectbox("标签", all_tags)
            
            # Build filter dict
            filters = {}
            if selected_type:
                filters["类型"] = selected_type
            if selected_name:
                filters["名称"] = selected_name
            if selected_tag:
                filters["标签"] = selected_tag
            
            # Show selected filters
            if filters:
                st.write("已选择的过滤条件:")
                for key, value in filters.items():
                    st.write(f"- {key}: {value}")
            else:
                st.info("没有选择过滤条件，将使用所有文件回答问题。")
            
            # Chat history
            if "rag_chat_history" not in st.session_state:
                st.session_state["rag_chat_history"] = []
            
            # Display chat history
            st.subheader("聊天记录")
            for message in st.session_state["rag_chat_history"]:
                if message["role"] == "user":
                    st.markdown(f"**您:** {message['content']}")
                else:
                    st.markdown(f"**AI:** {message['content']}")
            
            # Chat input
            st.subheader("与知识库对话")
            with st.form("rag_chat_form", clear_on_submit=True):
                chat_input = st.text_input("输入您的问题：")
                submitted = st.form_submit_button("发送")
                
                if submitted and chat_input:
                    # Add user message to chat history
                    st.session_state["rag_chat_history"].append({"role": "user", "content": chat_input})
                    
                    # Generate response based on the local knowledge base (with filters)
                    response = chat_with_local_facts(chat_input, filters)
                    
                    # Add AI response to chat history
                    st.session_state["rag_chat_history"].append({"role": "assistant", "content": response})
                    
                    # Rerun to update the UI
                    st.rerun()
            
            # Clear chat history button
            if st.button("清空聊天记录", key="clear_rag_chat"):
                st.session_state["rag_chat_history"] = []
                st.success("聊天记录已清空！")
                st.rerun()

# ----------------------- Tab 3: Direct Chat -----------------------
with tabs[2]:
    st.header("直接聊天")
    st.write("基于 Qwen、大模型、本地信息库以及 Deepseek 模型，您可以直接与 AI 进行对话。")
    
    # 选择聊天模式：包含 Qwen、本地知识库 (RAG) 和 Deepseek 聊天选项
    chat_mode = st.radio("选择聊天模式", ("Qwen聊天", "Deepseek聊天"))
    
    # Chat input form (clears on submit)
    with st.form("chat_form", clear_on_submit=True):
        chat_input = st.text_input("输入您的消息：")
        submitted = st.form_submit_button("发送")
        if submitted and chat_input:
            st.session_state["chat_history"].append({"role": "user", "content": chat_input})
            if chat_mode == "Qwen聊天":
                reply = chat_with_qwen(chat_input)
            elif chat_mode == "Deepseek聊天":
                reply = chat_with_deepseek(chat_input)
            st.session_state["chat_history"].append({"role": "assistant", "content": reply})
    
    st.markdown("### 聊天记录")
    for message in st.session_state["chat_history"]:
        if message["role"] == "user":
            st.markdown(f"**您:** {message['content']}")
        else:
            st.markdown(f"**AI:** {message['content']}")
    
    # Clear chat history button
    if st.button("清空聊天记录", key="clear_direct_chat"):
        st.session_state["chat_history"] = []
        st.success("聊天记录已清空！")